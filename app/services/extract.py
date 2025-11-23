# app/services/extract.py
from pypdf import PdfReader
from pypdf.errors import PdfStreamError
from io import BytesIO
import os
import zipfile  # NOTE: DOCX/PPTX(Office Open XML)는 zip 포맷이기 때문에 추가
from typing import List, Optional

# NOTE: 아래 두 라이브러리는 DOCX/PPTX 지원용
# - python-docx
# - python-pptx
# 를 설치해야 동작함 (pip install python-docx python-pptx)
try:
    from docx import Document as DocxDocument
except ImportError:  # 런타임에 없으면 None으로 두고, 사용 시 에러 던지게 함
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def _extract_pdf_pages_from_bytes(data: bytes) -> List[str]:
    """
    NOTE: 기존 PDF 추출 로직을 그대로 함수로 분리.
    입력: PDF 바이너리
    출력: 페이지별 텍스트 리스트
    """
    try:
        reader = PdfReader(BytesIO(data))
        return [(p.extract_text() or "").strip() for p in reader.pages]
    except PdfStreamError:
        # pdfminer.six 폴백
        from pdfminer.high_level import extract_text  # type: ignore
        txt = extract_text(BytesIO(data))
        return [t.strip() for t in txt.split("\f") if t.strip()]


def _detect_office_type_from_zip(data: bytes) -> Optional[str]:
    """
    NOTE: DOCX/PPTX는 모두 ZIP 기반 포맷이므로,
    Zip 안에 포함된 폴더 구조로 타입을 판별한다.
    - word/ 가 있으면 DOCX
    - ppt/ 가 있으면 PPTX
    그 외는 None
    """
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            names = zf.namelist()
    except zipfile.BadZipFile:
        return None

    has_word = any(name.startswith("word/") for name in names)
    has_ppt = any(name.startswith("ppt/") for name in names)

    if has_word:
        return "docx"
    if has_ppt:
        return "pptx"
    return None


def _extract_docx_pages_from_bytes(data: bytes) -> List[str]:
    """
    NOTE: DOCX에서 문단(paragraph) 단위로 텍스트를 읽어서
    페이지처럼 1,2,3... 번호를 붙여 반환한다.
    """
    if DocxDocument is None:
        raise RuntimeError("python-docx가 설치되어 있지 않습니다. `pip install python-docx` 필요.")

    doc = DocxDocument(BytesIO(data))
    pages: List[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        pages.append(text)
    return pages


def _extract_pptx_pages_from_bytes(data: bytes) -> List[str]:
    """
    NOTE: PPTX의 경우 슬라이드 하나를 하나의 'page'로 보고
    슬라이드 내 텍스트들을 합쳐서 반환한다.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. `pip install python-pptx` 필요.")

    prs = Presentation(BytesIO(data))
    pages: List[str] = []

    for slide in prs.slides:
        texts: List[str] = []
        for shape in slide.shapes:
            # shape마다 text 속성이 다를 수 있으므로 널리 커버
            text = getattr(shape, "text", None)
            if text:
                text = text.strip()
                if text:
                    texts.append(text)

        merged = "\n".join(texts).strip()
        if merged:
            pages.append(merged)

    return pages


def _extract_plain_text_pages_from_bytes(data: bytes) -> List[str]:
    """
    NOTE: TXT/MD 파일 등 '그냥 텍스트'로 볼 수 있는 경우.
    일단 전체를 UTF-8로 디코딩하고 하나의 page로 반환한다.
    (실제 페이지/섹션 분할은 chunker.py가 담당)
    """
    try:
        txt = data.decode("utf-8")
    except UnicodeDecodeError:
        # 필요하면 다른 인코딩을 시도할 수 있지만, MVP에선 utf-8 실패 시 에러
        raise ValueError("텍스트 파일을 utf-8로 디코딩할 수 없습니다.")

    txt = txt.strip()
    if not txt:
        return []
    return [txt]


def extract_text_pages(fp) -> list[str]:
    """
    NOTE: 기존 시그니처 유지 (indexer에서 그대로 사용 가능).
    변경점:
    - 더 이상 "PDF 아니면 무조건 에러"를 던지지 않고
    - 바이너리 내용을 보고 pdf / docx / pptx / txt 중 하나로 자동 판별해서 처리한다.
    """
    # 입력을 bytes로 정규화
    if isinstance(fp, (bytes, bytearray)):
        data = bytes(fp)
    elif hasattr(fp, "read"):  # file-like
        try:
            fp.seek(0)
        except Exception:
            pass
        data = fp.read()
    elif isinstance(fp, (str, os.PathLike)):
        with open(fp, "rb") as f:
            data = f.read()
    else:
        raise TypeError("Unsupported input type for extract_text_pages")

    # 빈 입력 방어
    if not data or len(data) < 4:
        raise ValueError("Input is too short to determine file type.")

    # 1) PDF 헤더 검사
    if data.lstrip().startswith(b"%PDF-"):
        # 기존 로직 그대로 사용
        return _extract_pdf_pages_from_bytes(data)

    # 2) ZIP 기반 Office 포맷(DOCX/PPTX) 검사
    if data.startswith(b"PK\x03\x04"):  # 일반적인 ZIP 시그니처
        office_type = _detect_office_type_from_zip(data)
        if office_type == "docx":
            return _extract_docx_pages_from_bytes(data)
        elif office_type == "pptx":
            return _extract_pptx_pages_from_bytes(data)
        # ZIP이지만 Office가 아닌 경우 → 아래 plain text 시도로 폴백 (또는 에러 던져도 됨)

    # 3) 나머지는 TXT/MD 등 '평문 텍스트'로 취급
    return _extract_plain_text_pages_from_bytes(data)
